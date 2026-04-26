#!/usr/bin/env python3
"""
Alche "Imperfect Longevity" — 4-Axis Positioning Deck Builder
Generates a 13-slide investor deck with an embedded positioning map.
"""

import os
import math
from io import BytesIO

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patheffects as pe
import numpy as np
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Design tokens ──────────────────────────────────────────────────────────────

BG       = RGBColor(0x0A, 0x0A, 0x0A)
GOLD     = RGBColor(0xC8, 0xA4, 0x4E)
CREAM    = RGBColor(0xF5, 0xF0, 0xEB)
MUTED    = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
CARD_BG  = RGBColor(0x14, 0x14, 0x14)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
MARGIN  = Inches(0.7)
GOLD_BAR_H = Emu(36576)  # 4px gold bar

FONT = "Arial"

# ── Brand data ─────────────────────────────────────────────────────────────────

def define_brand_data():
    """All 30 brands + ALCHE with 4-axis scores."""
    return [
        # Longevity Purists
        {"name": "Blueprint",        "cat": "Longevity Purist",  "x": 9.5, "y": 1.0, "size": 1.5, "color": 2},
        {"name": "Function Health",  "cat": "Longevity Purist",  "x": 8.5, "y": 3.0, "size": 3.0, "color": 3},
        {"name": "Ultrahuman",       "cat": "Longevity Purist",  "x": 8.0, "y": 2.5, "size": 2.0, "color": 3},
        {"name": "Fountain Life",    "cat": "Longevity Purist",  "x": 8.5, "y": 2.0, "size": 2.5, "color": 2},
        {"name": "NOVOS Labs",       "cat": "Longevity Purist",  "x": 7.5, "y": 3.5, "size": 3.0, "color": 2},
        {"name": "Superpower",       "cat": "Longevity Purist",  "x": 7.0, "y": 3.0, "size": 2.5, "color": 2},
        # Wearable/Data
        {"name": "WHOOP",            "cat": "Wearable/Data",     "x": 7.5, "y": 2.5, "size": 2.0, "color": 3},
        {"name": "Oura",             "cat": "Wearable/Data",     "x": 7.0, "y": 4.0, "size": 4.5, "color": 3},
        {"name": "Eight Sleep",      "cat": "Wearable/Data",     "x": 7.0, "y": 3.5, "size": 3.0, "color": 3},
        {"name": "Levels",           "cat": "Wearable/Data",     "x": 7.5, "y": 3.0, "size": 2.5, "color": 3},
        {"name": "Apple Health",     "cat": "Wearable/Data",     "x": 5.0, "y": 6.0, "size": 7.0, "color": 1},
        # Clinical/Fitness
        {"name": "Parsley Health",   "cat": "Clinical/Fitness",  "x": 7.0, "y": 5.0, "size": 6.0, "color": 1},
        {"name": "Equinox",          "cat": "Clinical/Fitness",  "x": 4.0, "y": 4.5, "size": 4.0, "color": 3},
        {"name": "Life Time",        "cat": "Clinical/Fitness",  "x": 4.0, "y": 5.5, "size": 5.5, "color": 4},
        {"name": "Noom",             "cat": "Clinical/Fitness",  "x": 5.5, "y": 5.5, "size": 6.5, "color": 3},
        # Premium Lifestyle
        {"name": "Goop",             "cat": "Premium Lifestyle", "x": 2.5, "y": 7.0, "size": 7.0, "color": 3},
        {"name": "Erewhon",          "cat": "Premium Lifestyle", "x": 1.5, "y": 7.5, "size": 5.0, "color": 3},
        {"name": "Aman Wellness",    "cat": "Premium Lifestyle", "x": 3.0, "y": 7.0, "size": 4.0, "color": 3},
        {"name": "Six Senses",       "cat": "Premium Lifestyle", "x": 3.0, "y": 6.5, "size": 4.0, "color": 2},
        # Product Brands
        {"name": "AG1",              "cat": "Product Brand",     "x": 4.0, "y": 5.0, "size": 5.0, "color": 3},
        {"name": "Seed",             "cat": "Product Brand",     "x": 6.0, "y": 4.0, "size": 4.5, "color": 2},
        {"name": "Thorne",           "cat": "Product Brand",     "x": 6.5, "y": 2.5, "size": 3.0, "color": 1},
        {"name": "Moon Juice",       "cat": "Product Brand",     "x": 2.0, "y": 6.5, "size": 6.0, "color": 2},
        {"name": "Momentous",        "cat": "Product Brand",     "x": 6.0, "y": 3.0, "size": 2.5, "color": 1},
        # Anti-Establishment
        {"name": "Liquid Death",     "cat": "Anti-Establishment","x": 0.5, "y": 9.0, "size": 7.5, "color": 1},
        {"name": "MUD\\WTR",         "cat": "Anti-Establishment","x": 2.0, "y": 8.5, "size": 6.5, "color": 1},
        # Berlin / Culture
        {"name": "ANTI Berlin",      "cat": "Berlin/Culture",    "x": 3.0, "y": 8.0, "size": 5.5, "color": 2},
        {"name": "Daluma Berlin",    "cat": "Berlin/Culture",    "x": 2.5, "y": 7.0, "size": 5.0, "color": 2},
        # Longevity Skincare
        {"name": "OneSkin",          "cat": "Longevity Skincare","x": 7.0, "y": 5.0, "size": 6.5, "color": 1},
        {"name": "Timeline",         "cat": "Longevity Skincare","x": 7.5, "y": 3.5, "size": 4.0, "color": 1},
        # ALCHE
        {"name": "ALCHE",            "cat": "ALCHE",             "x": 8.0, "y": 8.5, "size": 9.0, "color": 6},
    ]


# ── Chart generation ───────────────────────────────────────────────────────────

def generate_positioning_map(brands, path):
    """Create the 4-axis positioning scatter chart."""
    fig, ax = plt.subplots(figsize=(12, 7.5))
    fig.patch.set_facecolor("#0A0A0A")
    ax.set_facecolor("#0A0A0A")

    # Quadrant dividers
    ax.axhline(y=5.0, color="#333333", linewidth=0.8, linestyle="-", zorder=1)
    ax.axvline(x=5.0, color="#333333", linewidth=0.8, linestyle="-", zorder=1)

    # Dashed gold rectangle for top-right quadrant (the "white space")
    from matplotlib.patches import FancyBboxPatch, Rectangle
    rect = Rectangle((5.0, 5.0), 5.5, 5.5, linewidth=1.8, edgecolor="#C8A44E",
                      facecolor="#C8A44E", alpha=0.04, linestyle="--", zorder=1)
    ax.add_patch(rect)

    # Quadrant labels
    label_props = dict(fontsize=9, fontfamily="Arial", alpha=0.35, ha="center")
    ax.text(2.5, 0.4, "UNDIFFERENTIATED", color="#888888", **label_props)
    ax.text(7.5, 0.4, "OPTIMIZATION-FIRST", color="#888888", **label_props)
    ax.text(2.5, 9.8, "LIFESTYLE WELLNESS", color="#888888", **label_props)
    ax.text(7.5, 9.8, "UNOCCUPIED", color="#C8A44E", fontsize=10,
            fontfamily="Arial", alpha=0.6, ha="center", fontweight="bold")

    # Color mapping: integration score 1-6 → color gradient (muted red → gold)
    def score_to_color(score):
        t = (score - 1) / 5.0
        r = int(140 + t * (200 - 140))
        g = int(50 + t * (164 - 50))
        b = int(50 + t * (78 - 50))
        return (r / 255, g / 255, b / 255)

    # Plot non-ALCHE brands
    for b in brands:
        if b["name"] == "ALCHE":
            continue
        s = (b["size"] * 18) ** 1.6  # scale bubble area
        c = score_to_color(b["color"])
        ax.scatter(b["x"], b["y"], s=s, c=[c], alpha=0.65, edgecolors="none", zorder=3)
        # Label
        offset_x = 0.18
        offset_y = -0.28
        # Adjust labels for crowded areas
        if b["name"] in ("NOVOS Labs", "Timeline"):
            offset_y = 0.28
        if b["name"] == "Levels":
            offset_x = 0.25
            offset_y = 0.2
        if b["name"] == "Eight Sleep":
            offset_y = -0.35
        if b["name"] == "Superpower":
            offset_y = 0.25
        if b["name"] == "Daluma Berlin":
            offset_x = -0.3
            offset_y = -0.5
        if b["name"] == "OneSkin":
            offset_x = -0.2
            offset_y = 0.3
        if b["name"] == "Parsley Health":
            offset_x = -0.3
            offset_y = -0.3
        if b["name"] == "Aman Wellness":
            offset_x = 0.3
            offset_y = -0.45
        if b["name"] == "Goop":
            offset_x = -0.4
            offset_y = -0.05
        if b["name"] == "Six Senses":
            offset_x = 0.3
            offset_y = -0.1
        ax.text(b["x"] + offset_x, b["y"] + offset_y, b["name"],
                fontsize=6.5, color="#AAAAAA", fontfamily="Arial", alpha=0.8, zorder=4)

    # Plot ALCHE with glow
    alche = [b for b in brands if b["name"] == "ALCHE"][0]
    # Outer glow
    ax.scatter(alche["x"], alche["y"], s=2800, c=["#C8A44E"], alpha=0.12, edgecolors="none", zorder=5)
    ax.scatter(alche["x"], alche["y"], s=1600, c=["#C8A44E"], alpha=0.20, edgecolors="none", zorder=5)
    # Main bubble
    ax.scatter(alche["x"], alche["y"], s=900, c=["#C8A44E"], alpha=0.9,
               edgecolors="#F5F0EB", linewidth=1.5, zorder=6)
    ax.text(alche["x"] + 0.35, alche["y"] + 0.05, "ALCHE", fontsize=11, color="#F5F0EB",
            fontfamily="Arial", fontweight="bold", zorder=7,
            path_effects=[pe.withStroke(linewidth=3, foreground="#0A0A0A")])

    # Axes
    ax.set_xlim(-0.3, 10.5)
    ax.set_ylim(-0.3, 10.5)
    ax.set_xlabel("Lifestyle / Aesthetic  ←――――――――――→  Clinical / Data-Driven",
                  fontsize=9, color="#999999", fontfamily="Arial", labelpad=10)
    ax.set_ylabel("Ascetic / Optimization  ←――――――――――→  Human / Joyful",
                  fontsize=9, color="#999999", fontfamily="Arial", labelpad=10)

    ax.tick_params(colors="#555555", labelsize=7)
    for spine in ax.spines.values():
        spine.set_color("#333333")
        spine.set_linewidth(0.5)

    ax.set_xticks(range(0, 11))
    ax.set_yticks(range(0, 11))
    ax.grid(True, alpha=0.08, color="#555555", linewidth=0.5)

    # Legend — bubble size
    legend_x_start = 0.02
    legend_y = 0.04
    for val, label in [(3, "Low inclusion"), (6, "Medium"), (9, "High inclusion")]:
        s = (val * 18) ** 1.6
        ax.scatter([], [], s=s, c=["#999999"], alpha=0.5, edgecolors="none", label=label)

    leg = ax.legend(loc="lower left", frameon=False, fontsize=7,
                    labelcolor="#999999", handletextpad=1.2, borderpad=1.0,
                    title="Bubble size = Inclusivity", title_fontproperties={"size": 7.5})
    leg.get_title().set_color("#999999")

    # Color bar legend for integration score
    from matplotlib.colors import LinearSegmentedColormap
    cmap_colors = [score_to_color(i) for i in range(1, 7)]
    cmap = LinearSegmentedColormap.from_list("integration", cmap_colors, N=6)
    sm = plt.cm.ScalarMappable(cmap=cmap, norm=plt.Normalize(vmin=1, vmax=6))
    sm.set_array([])
    cbar_ax = fig.add_axes([0.78, 0.08, 0.15, 0.015])
    cbar = fig.colorbar(sm, cax=cbar_ax, orientation="horizontal", ticks=[1, 2, 3, 4, 5, 6])
    cbar.ax.tick_params(labelsize=6, colors="#999999", length=2)
    cbar.outline.set_edgecolor("#333333")
    cbar.outline.set_linewidth(0.5)
    cbar.set_label("Integration pillars (1–6)", fontsize=7, color="#999999",
                   fontfamily="Arial", labelpad=4)

    fig.subplots_adjust(left=0.08, right=0.95, top=0.95, bottom=0.12)
    fig.savefig(path, dpi=300, facecolor="#0A0A0A", bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)
    print(f"  Chart saved: {path}")


# ── Slide helpers ──────────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG):
    """Set solid background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gold_bar(slide):
    """4px gold accent bar at top of slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, GOLD_BAR_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=CREAM, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT, line_spacing=None):
    """Add a text box and return it."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.bold = bold
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox

def add_section_number(slide, number):
    """Section number in bottom-right corner."""
    add_text_box(slide, Inches(8.5), Inches(4.8), Inches(0.8), Inches(0.5),
                 f"{number:02d}", font_size=24, color=DARK_GRAY,
                 alignment=PP_ALIGN.RIGHT, bold=False)

def add_multi_run_text(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT, line_spacing=None):
    """Add a text box with multiple styled runs. runs = list of (text, size, color, bold)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    for i, (text, size, color, bold) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT
        run.font.bold = bold
    return txBox

def add_card(slide, left, top, width, height, brand_name, scores_text, descriptor,
             bg_color=CARD_BG, border_color=None):
    """Add a brand card rectangle with text inside."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    # Brand name
    add_text_box(slide, left + Inches(0.15), top + Inches(0.12),
                 width - Inches(0.3), Inches(0.3),
                 brand_name, font_size=14, color=CREAM, bold=True)
    # Scores
    add_text_box(slide, left + Inches(0.15), top + Inches(0.42),
                 width - Inches(0.3), Inches(0.25),
                 scores_text, font_size=10, color=GOLD, bold=False)
    # Descriptor
    add_text_box(slide, left + Inches(0.15), top + Inches(0.68),
                 width - Inches(0.3), Inches(0.5),
                 descriptor, font_size=9, color=MUTED, bold=False)


# ── Slide builders ─────────────────────────────────────────────────────────────

def build_slide_01_title(prs):
    """Slide 1 — Title hero."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    add_gold_bar(slide)

    # ALCHE wordmark
    add_text_box(slide, MARGIN, Inches(0.5), Inches(3), Inches(0.4),
                 "ALCHE", font_size=16, color=GOLD, bold=True)

    # Headline — split into two color runs
    add_multi_run_text(slide, MARGIN, Inches(1.1), Inches(8), Inches(1.6), [
        ("No longevity brand occupies the intersection\nof deep science ", 40, CREAM, True),
        ("and human joy", 40, GOLD, True),
    ], line_spacing=48)

    # Subtitle
    add_text_box(slide, MARGIN, Inches(2.8), Inches(7.5), Inches(0.4),
                 "Evidence from 30 brands across the $6.8T wellness economy",
                 font_size=13, color=LIGHT_GRAY, bold=False)

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  MARGIN, Inches(3.45), Inches(7.5), Emu(9144))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    line.line.fill.background()

    # Stat row
    stats = [
        ("30", "BRANDS ANALYZED"),
        ("4", "POSITIONING AXES"),
        ("0", "IN THE WHITE SPACE"),
    ]
    for i, (val, label) in enumerate(stats):
        x = MARGIN + Inches(i * 2.5)
        add_text_box(slide, x, Inches(3.65), Inches(1.5), Inches(0.5),
                     val, font_size=36, color=GOLD, bold=True)
        add_text_box(slide, x, Inches(4.2), Inches(1.8), Inches(0.25),
                     label, font_size=8, color=MUTED, bold=False)

    # Bottom credit
    add_text_box(slide, MARGIN, Inches(5.05), Inches(5), Inches(0.25),
                 "Competitive Positioning Analysis  ·  February 2026",
                 font_size=8, color=DARK_GRAY, bold=False)


def build_slide_02_retention(prs):
    """Slide 2 — The Retention Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)
    add_section_number(slide, 1)

    # Section label
    add_text_box(slide, MARGIN, Inches(0.35), Inches(4), Inches(0.25),
                 "01 — THE RETENTION PROBLEM", font_size=9, color=MUTED, bold=False)

    # Headline
    add_text_box(slide, MARGIN, Inches(0.65), Inches(8.2), Inches(0.7),
                 "90% of wearable users abandon their device within 12 months",
                 font_size=32, color=CREAM, bold=True, line_spacing=38)

    # Large stat
    add_text_box(slide, MARGIN, Inches(1.7), Inches(3.5), Inches(1.8),
                 "90%", font_size=96, color=GOLD, bold=True)

    # Supporting label under stat
    add_text_box(slide, MARGIN, Inches(3.3), Inches(3.5), Inches(0.3),
                 "wearable abandonment rate", font_size=12, color=MUTED, bold=False)

    # Right column — supporting data
    right_x = Inches(5.2)
    add_text_box(slide, right_x, Inches(1.9), Inches(4), Inches(0.5),
                 "3.4%", font_size=48, color=GOLD, bold=True)
    add_text_box(slide, right_x, Inches(2.55), Inches(4), Inches(0.3),
                 "health app retention at day 30", font_size=12, color=MUTED, bold=False)

    # Divider
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  right_x, Inches(3.1), Inches(3.8), Emu(9144))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    line.line.fill.background()

    # Takeaway
    add_text_box(slide, right_x, Inches(3.35), Inches(4), Inches(0.8),
                 "The current model optimizes for measurement.\nUsers leave anyway.",
                 font_size=13, color=CREAM, bold=False, line_spacing=19)

    # Source
    add_text_box(slide, MARGIN, Inches(5.05), Inches(6), Inches(0.25),
                 "Sources: Industry research data on wearable retention and health app engagement",
                 font_size=7, color=DARK_GRAY, bold=False)


def build_slide_03_backlash(prs):
    """Slide 3 — The Backlash."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)
    add_section_number(slide, 2)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(4), Inches(0.25),
                 "02 — THE BACKLASH", font_size=9, color=MUTED, bold=False)

    add_text_box(slide, MARGIN, Inches(0.65), Inches(8.5), Inches(0.7),
                 "The Global Wellness Summit named the over-optimization backlash its #2 trend for 2026",
                 font_size=28, color=CREAM, bold=True, line_spacing=34)

    # Left column — pull quote
    left_x = MARGIN
    # Quote mark
    add_text_box(slide, left_x, Inches(1.65), Inches(0.5), Inches(0.5),
                 "\u201C", font_size=48, color=GOLD, bold=True)

    add_text_box(slide, left_x + Inches(0.05), Inches(2.05), Inches(4.3), Inches(1.3),
                 "Wellness experiences will embrace what humans actually are: imperfect, emotional, relational and sensory.",
                 font_size=15, color=CREAM, bold=False, line_spacing=22)

    add_text_box(slide, left_x + Inches(0.05), Inches(3.35), Inches(4), Inches(0.25),
                 "— Global Wellness Summit, Future of Wellness 2026",
                 font_size=9, color=MUTED, bold=False)

    # Right column — evidence points
    right_x = Inches(5.5)
    points = [
        ("Trend #2 of 2026:", "Rehumanizing wellness — moving from metrics-obsessed to experience-centered"),
        ("Consumer signal:", "The data-overload fatigue driving users away from wearables is cultural, not technical"),
        ("Market implication:", "Brands that combine scientific rigor with emotional resonance will capture the backlash demand"),
    ]
    for i, (head, body) in enumerate(points):
        y = Inches(1.7) + Inches(i * 1.0)
        add_text_box(slide, right_x, y, Inches(3.8), Inches(0.25),
                     head, font_size=11, color=GOLD, bold=True)
        add_text_box(slide, right_x, y + Inches(0.25), Inches(3.8), Inches(0.55),
                     body, font_size=10, color=LIGHT_GRAY, bold=False, line_spacing=14)

    add_text_box(slide, MARGIN, Inches(5.05), Inches(7), Inches(0.25),
                 "Sources: GWS Future of Wellness 2026, PR Newswire, Global Wellness Institute",
                 font_size=7, color=DARK_GRAY, bold=False)


def build_slide_04_gender_gap(prs):
    """Slide 4 — The Gender Gap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)
    add_section_number(slide, 3)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(4), Inches(0.25),
                 "03 — THE GENDER GAP", font_size=9, color=MUTED, bold=False)

    add_text_box(slide, MARGIN, Inches(0.65), Inches(8.5), Inches(0.7),
                 "Women\u2019s health receives 6% of private healthcare investment despite being 51% of the population",
                 font_size=28, color=CREAM, bold=True, line_spacing=34)

    # Left column — stat callouts
    left_x = MARGIN
    stat_data = [
        ("6%", "of private healthcare investment\ngoes to women's health", "WEF / Cosmeticsdesign-europe"),
        ("$40B", "menopause market by 2030\n(8× growth from 2023)", "BCG"),
        ("$100M", "committed by Melinda French Gates\nto women's health R&D", "Pivotal Ventures"),
    ]
    for i, (val, desc, src) in enumerate(stat_data):
        y = Inches(1.6) + Inches(i * 1.05)
        add_text_box(slide, left_x, y, Inches(1.8), Inches(0.55),
                     val, font_size=36, color=GOLD, bold=True)
        add_text_box(slide, left_x + Inches(1.9), y + Inches(0.05), Inches(2.8), Inches(0.55),
                     desc, font_size=10, color=LIGHT_GRAY, bold=False, line_spacing=14)
        add_text_box(slide, left_x + Inches(1.9), y + Inches(0.6), Inches(2.8), Inches(0.2),
                     f"Source: {src}", font_size=7, color=DARK_GRAY, bold=False)

    # Right column — context
    right_x = Inches(5.5)

    # Quote
    add_text_box(slide, right_x, Inches(1.7), Inches(0.4), Inches(0.4),
                 "\u201C", font_size=36, color=GOLD, bold=True)
    add_text_box(slide, right_x + Inches(0.05), Inches(2.0), Inches(3.8), Inches(1.0),
                 "The booming longevity market, like medicine decades before it, is tacitly male.",
                 font_size=14, color=CREAM, bold=False, line_spacing=20)
    add_text_box(slide, right_x + Inches(0.05), Inches(2.9), Inches(3.8), Inches(0.25),
                 "— Global Wellness Summit 2026", font_size=9, color=MUTED, bold=False)

    # Annotation
    add_text_box(slide, right_x, Inches(3.5), Inches(3.8), Inches(1.0),
                 "This creates a structural opportunity: the first longevity platform to center women\u2019s health captures an underserved majority.",
                 font_size=11, color=LIGHT_GRAY, bold=False, line_spacing=16)


def build_slide_05_positioning_map(prs, chart_path):
    """Slide 5 — THE POSITIONING MAP (centerpiece)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)
    add_section_number(slide, 4)

    add_text_box(slide, MARGIN, Inches(0.2), Inches(4), Inches(0.2),
                 "04 — POSITIONING MAP", font_size=9, color=MUTED, bold=False)

    add_text_box(slide, MARGIN, Inches(0.42), Inches(8.5), Inches(0.45),
                 "30 brands plotted on 4 axes \u2014 the top-right quadrant is empty",
                 font_size=22, color=CREAM, bold=True)

    # Embed chart image — near full bleed
    chart_left = Inches(0.2)
    chart_top = Inches(0.92)
    chart_width = Inches(9.6)
    # Calculate height from image aspect ratio
    img = Image.open(chart_path)
    aspect = img.height / img.width
    chart_height_val = 9.6 * aspect
    # Cap to fit slide
    max_h = 4.5
    if chart_height_val > max_h:
        chart_width = Inches(max_h / aspect)
        chart_height = Inches(max_h)
        chart_left = Inches((10 - max_h / aspect) / 2)
    else:
        chart_height = Inches(chart_height_val)

    slide.shapes.add_picture(chart_path, chart_left, chart_top, chart_width, chart_height)


def build_slide_06_optimization_quadrant(prs):
    """Slide 6 — Optimization-First Quadrant."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)
    add_section_number(slide, 5)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(4), Inches(0.25),
                 "05 — OPTIMIZATION-FIRST QUADRANT", font_size=9, color=MUTED, bold=False)

    add_text_box(slide, MARGIN, Inches(0.65), Inches(8.5), Inches(0.7),
                 "High-science brands default to ascetic positioning and male-oriented design",
                 font_size=28, color=CREAM, bold=True, line_spacing=34)

    # 2×2 card grid
    cards = [
        ("Blueprint", "Science: 9.5  ·  Humanity: 1.0  ·  Inclusion: 1.5",
         "Bryan Johnson\u2019s protocol-driven approach. Highest science score, lowest humanity score in dataset."),
        ("Function Health", "Science: 8.5  ·  Humanity: 3.0  ·  Inclusion: 3.0",
         "100+ biomarker testing platform. Clinical positioning with limited lifestyle integration."),
        ("WHOOP", "Science: 7.5  ·  Humanity: 2.5  ·  Inclusion: 2.0",
         "Performance-first wearable. Core audience: male athletes and biohackers."),
        ("Ultrahuman", "Science: 8.0  ·  Humanity: 2.5  ·  Inclusion: 2.0",
         "Metabolic tracking via CGM. Tech-forward brand language, narrow demographic appeal."),
    ]
    card_w = Inches(4.05)
    card_h = Inches(1.45)
    gap = Inches(0.25)
    start_x = MARGIN
    start_y = Inches(1.6)

    for i, (name, scores, desc) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)
        add_card(slide, x, y, card_w, card_h, name, scores, desc)


def build_slide_07_lifestyle_quadrant(prs):
    """Slide 7 — Lifestyle Wellness Quadrant."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)
    add_section_number(slide, 6)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(4), Inches(0.25),
                 "06 — LIFESTYLE WELLNESS QUADRANT", font_size=9, color=MUTED, bold=False)

    add_text_box(slide, MARGIN, Inches(0.65), Inches(8.5), Inches(0.7),
                 "High-humanity brands lack clinical depth or longevity specificity",
                 font_size=28, color=CREAM, bold=True, line_spacing=34)

    cards = [
        ("Goop", "Science: 2.5  ·  Humanity: 7.0  ·  Inclusion: 7.0",
         "Lifestyle-first wellness brand. Strong cultural resonance but limited scientific credibility."),
        ("Erewhon", "Science: 1.5  ·  Humanity: 7.5  ·  Inclusion: 5.0",
         "Premium grocery as lifestyle brand. High aspiration, minimal health-outcome depth."),
        ("Moon Juice", "Science: 2.0  ·  Humanity: 6.5  ·  Inclusion: 6.0",
         "Adaptogen-forward product brand. Aesthetic strength without clinical grounding."),
        ("Aman Wellness", "Science: 3.0  ·  Humanity: 7.0  ·  Inclusion: 4.0",
         "Ultra-luxury wellness hospitality. Beautiful experience, low science integration."),
    ]
    card_w = Inches(4.05)
    card_h = Inches(1.45)
    gap = Inches(0.25)
    start_x = MARGIN
    start_y = Inches(1.6)

    for i, (name, scores, desc) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)
        add_card(slide, x, y, card_w, card_h, name, scores, desc)


def build_slide_08_anti_establishment(prs):
    """Slide 8 — Anti-Establishment Precedent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)
    add_section_number(slide, 7)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(4), Inches(0.25),
                 "07 — CATEGORY INVERSION PRECEDENT", font_size=9, color=MUTED, bold=False)

    add_text_box(slide, MARGIN, Inches(0.65), Inches(8.5), Inches(0.7),
                 "Liquid Death reached $1.4B valuation by inverting category conventions in water",
                 font_size=28, color=CREAM, bold=True, line_spacing=34)

    # Left case study
    left_x = MARGIN
    case_w = Inches(3.9)

    # Liquid Death card
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left_x, Inches(1.65), case_w, Inches(2.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()

    add_text_box(slide, left_x + Inches(0.2), Inches(1.8), Inches(3.5), Inches(0.35),
                 "Liquid Death", font_size=18, color=CREAM, bold=True)
    add_text_box(slide, left_x + Inches(0.2), Inches(2.15), Inches(3.5), Inches(0.25),
                 "Commoditized water → punk aesthetic", font_size=10, color=GOLD, bold=False)

    ld_stats = [
        ("$1.4B", "valuation"),
        ("$333M", "revenue (2024)"),
    ]
    for j, (val, lab) in enumerate(ld_stats):
        sy = Inches(2.6) + Inches(j * 0.65)
        add_text_box(slide, left_x + Inches(0.2), sy, Inches(1.5), Inches(0.4),
                     val, font_size=28, color=GOLD, bold=True)
        add_text_box(slide, left_x + Inches(1.8), sy + Inches(0.08), Inches(2), Inches(0.3),
                     lab, font_size=10, color=MUTED, bold=False)

    add_text_box(slide, left_x + Inches(0.2), Inches(3.9), Inches(3.5), Inches(0.3),
                 "Sources: Fueler, Sacra, Entrepreneur", font_size=7, color=DARK_GRAY, bold=False)

    # Right case study — MUD\WTR
    right_x = Inches(5.2)
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    right_x, Inches(1.65), case_w, Inches(2.8))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = CARD_BG
    shape2.line.fill.background()

    add_text_box(slide, right_x + Inches(0.2), Inches(1.8), Inches(3.5), Inches(0.35),
                 "MUD\\WTR", font_size=18, color=CREAM, bold=True)
    add_text_box(slide, right_x + Inches(0.2), Inches(2.15), Inches(3.5), Inches(0.25),
                 "Anti-hustle-culture coffee alternative", font_size=10, color=GOLD, bold=False)

    mw_stats = [
        ("$47M", "revenue (2024 est.)"),
    ]
    for j, (val, lab) in enumerate(mw_stats):
        sy = Inches(2.6) + Inches(j * 0.65)
        add_text_box(slide, right_x + Inches(0.2), sy, Inches(1.5), Inches(0.4),
                     val, font_size=28, color=GOLD, bold=True)
        add_text_box(slide, right_x + Inches(1.8), sy + Inches(0.08), Inches(2), Inches(0.3),
                     lab, font_size=10, color=MUTED, bold=False)

    add_text_box(slide, right_x + Inches(0.2), Inches(3.9), Inches(3.5), Inches(0.3),
                 "Sources: ecdb, Owler", font_size=7, color=DARK_GRAY, bold=False)

    # Bottom takeaway
    add_text_box(slide, MARGIN, Inches(4.7), Inches(8), Inches(0.5),
                 "Neither operates in longevity. The archetype is proven but the category is untouched.",
                 font_size=13, color=CREAM, bold=False)


def build_slide_09_proposed_position(prs):
    """Slide 9 — The Proposed Position."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)
    add_section_number(slide, 8)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(4), Inches(0.25),
                 "08 — THE PROPOSED POSITION", font_size=9, color=MUTED, bold=False)

    add_text_box(slide, MARGIN, Inches(0.65), Inches(8.5), Inches(0.7),
                 "Alche sits at the intersection: clinical-grade science, inclusive design, joyful brand",
                 font_size=28, color=CREAM, bold=True, line_spacing=34)

    # Three equal columns
    col_w = Inches(2.65)
    gap = Inches(0.2)
    col_h = Inches(3.0)
    start_y = Inches(1.65)

    columns = [
        ("Clinical Science", "8.0 / 10", [
            "Bloodwork & biomarker panels",
            "AI-personalized protocols",
            "Evidence-based supplementation",
            "Longitudinal health tracking",
        ]),
        ("Human Brand", "8.5 / 10", [
            "Permission-giving, not punishing",
            "Lifestyle integration over restriction",
            "Joyful aesthetics & cultural edge",
            "Community over competition",
        ]),
        ("Inclusive Design", "9.0 / 10", [
            "Women-first approach",
            "Hormonal health integration",
            "Diverse bodies & life stages",
            "Accessible entry points",
        ]),
    ]

    for i, (title, score, bullets) in enumerate(columns):
        x = MARGIN + i * (col_w + gap)

        # Column card
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, start_y, col_w, col_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()

        # Title
        add_text_box(slide, x + Inches(0.15), start_y + Inches(0.12),
                     col_w - Inches(0.3), Inches(0.3),
                     title, font_size=15, color=GOLD, bold=True)

        # Score
        add_text_box(slide, x + Inches(0.15), start_y + Inches(0.42),
                     col_w - Inches(0.3), Inches(0.25),
                     score, font_size=11, color=MUTED, bold=False)

        # Gold accent line
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        x + Inches(0.15), start_y + Inches(0.72),
                                        Inches(0.8), Emu(13716))
        accent.fill.solid()
        accent.fill.fore_color.rgb = GOLD
        accent.line.fill.background()

        # Bullet points
        for j, bullet in enumerate(bullets):
            add_text_box(slide, x + Inches(0.15), start_y + Inches(0.9 + j * 0.45),
                         col_w - Inches(0.3), Inches(0.4),
                         f"\u2022  {bullet}", font_size=10, color=LIGHT_GRAY, bold=False)


def build_slide_10_integration_moat(prs):
    """Slide 10 — Integration Moat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)
    add_section_number(slide, 9)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(4), Inches(0.25),
                 "09 — INTEGRATION MOAT", font_size=9, color=MUTED, bold=False)

    add_text_box(slide, MARGIN, Inches(0.65), Inches(8.5), Inches(0.7),
                 "Alche scores 6/6 on integration pillars; the highest competitor scores 4/6",
                 font_size=28, color=CREAM, bold=True, line_spacing=34)

    # Pillars
    pillars = ["Physical\nSpace", "Products", "AI App", "Membership", "Community", "Integrated\nInfo System"]

    # Brands to compare
    brands_compare = [
        ("ALCHE",     [1, 1, 1, 1, 1, 1], GOLD),
        ("Life Time", [1, 1, 1, 1, 0, 0], RGBColor(0x88, 0x77, 0x55)),
        ("WHOOP",     [0, 1, 1, 1, 0, 0], RGBColor(0x66, 0x55, 0x44)),
        ("Equinox",   [1, 0, 1, 1, 0, 0], RGBColor(0x55, 0x44, 0x33)),
    ]

    # Draw horizontal bar chart
    bar_start_x = Inches(2.2)
    bar_w_unit = Inches(1.05)
    bar_h = Inches(0.5)
    gap_y = Inches(0.15)
    start_y = Inches(1.75)

    # Column headers (pillars)
    for j, pillar in enumerate(pillars):
        px = bar_start_x + j * bar_w_unit + Inches(0.15)
        add_text_box(slide, px, Inches(1.45), bar_w_unit - Inches(0.1), Inches(0.35),
                     pillar, font_size=7.5, color=MUTED, bold=False,
                     alignment=PP_ALIGN.CENTER)

    for i, (brand_name, scores, brand_color) in enumerate(brands_compare):
        y = start_y + i * (bar_h + gap_y)

        # Brand name
        add_text_box(slide, MARGIN, y + Inches(0.08), Inches(1.4), Inches(0.35),
                     brand_name, font_size=13,
                     color=CREAM if brand_name == "ALCHE" else LIGHT_GRAY,
                     bold=(brand_name == "ALCHE"))

        # Score
        total = sum(scores)
        add_text_box(slide, Inches(8.6), y + Inches(0.08), Inches(0.7), Inches(0.35),
                     f"{total}/6", font_size=13, color=brand_color, bold=True,
                     alignment=PP_ALIGN.RIGHT)

        # Pillar cells
        for j, has in enumerate(scores):
            cx = bar_start_x + j * bar_w_unit
            cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          cx, y, bar_w_unit - Inches(0.06), bar_h)
            cell.fill.solid()
            if has:
                cell.fill.fore_color.rgb = brand_color
                # Checkmark
                add_text_box(slide, cx + Inches(0.3), y + Inches(0.08),
                             Inches(0.4), Inches(0.3),
                             "\u2713", font_size=16,
                             color=RGBColor(0x0A, 0x0A, 0x0A) if brand_name == "ALCHE" else CREAM,
                             bold=True, alignment=PP_ALIGN.CENTER)
            else:
                cell.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
                add_text_box(slide, cx + Inches(0.3), y + Inches(0.08),
                             Inches(0.4), Inches(0.3),
                             "\u2014", font_size=14, color=RGBColor(0x44, 0x44, 0x44),
                             bold=False, alignment=PP_ALIGN.CENTER)
            cell.line.fill.background()

    # Source note
    add_text_box(slide, MARGIN, Inches(4.8), Inches(7), Inches(0.25),
                 "Source: Alche 6-pillar integration analysis of 30+ wellness & longevity companies",
                 font_size=7, color=DARK_GRAY, bold=False)


def build_slide_11_market_sizing(prs):
    """Slide 11 — Market Sizing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)
    add_section_number(slide, 10)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(4), Inches(0.25),
                 "10 — MARKET CONTEXT", font_size=9, color=MUTED, bold=False)

    add_text_box(slide, MARGIN, Inches(0.65), Inches(8.5), Inches(0.7),
                 "The wellness economy is $6.8T and growing to $10T by 2029",
                 font_size=28, color=CREAM, bold=True, line_spacing=34)

    # 2×3 stat grid
    stats = [
        ("$6.8T", "global wellness economy (2025)", "Global Wellness Institute"),
        ("$10T", "projected by 2029", "Global Wellness Institute"),
        ("$40B", "menopause market by 2030", "BCG"),
        ("57%", "prioritize aging well more\nthan 5 years ago", "NielsenIQ"),
        ("84%", "of consumers say wellness\nis a top priority", "McKinsey"),
        ("6%", "of private healthcare investment\ngoes to women\u2019s health", "WEF"),
    ]

    card_w = Inches(2.65)
    card_h = Inches(1.3)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = MARGIN
    start_y = Inches(1.6)

    for i, (val, desc, src) in enumerate(stats):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()

        add_text_box(slide, x + Inches(0.15), y + Inches(0.1),
                     card_w - Inches(0.3), Inches(0.45),
                     val, font_size=30, color=GOLD, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.55),
                     card_w - Inches(0.3), Inches(0.45),
                     desc, font_size=9, color=LIGHT_GRAY, bold=False, line_spacing=13)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.0),
                     card_w - Inches(0.3), Inches(0.2),
                     f"Source: {src}", font_size=7, color=DARK_GRAY, bold=False)


def build_slide_12_berlin(prs):
    """Slide 12 — Berlin Context."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)
    add_section_number(slide, 11)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(4), Inches(0.25),
                 "11 — BERLIN CONTEXT", font_size=9, color=MUTED, bold=False)

    add_text_box(slide, MARGIN, Inches(0.65), Inches(8.5), Inches(0.7),
                 "Berlin has an existing longevity scene and cultural DNA that aligns with the position",
                 font_size=28, color=CREAM, bold=True, line_spacing=34)

    # Three proof points as horizontal cards
    points = [
        ("ANTI Berlin",
         "Longevity + culture hybrid on Brunnenstraße. 20K Instagram followers. Single location proving the demand for science-meets-lifestyle in Berlin.",
         "IGNANT, antispaces.com"),
        ("Berlin Biohackers",
         "2,213 members in the Berlin Biohackers meetup group. Active community of self-experimenters already engaged with longevity practices.",
         "Meetup / Reddit research"),
        ("LIFE Summit 2026",
         "Berlin hosting a dedicated longevity summit. The city is becoming a European hub for longevity discourse and innovation.",
         "LIFE Summit"),
    ]

    start_y = Inches(1.6)
    card_w = Inches(8.6)
    card_h = Inches(1.05)
    gap = Inches(0.15)

    for i, (title, desc, src) in enumerate(points):
        y = start_y + i * (card_h + gap)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       MARGIN, y, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()

        # Number circle
        num_x = MARGIN + Inches(0.15)
        add_text_box(slide, num_x, y + Inches(0.08), Inches(0.35), Inches(0.35),
                     str(i + 1), font_size=16, color=GOLD, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, MARGIN + Inches(0.55), y + Inches(0.08),
                     Inches(3), Inches(0.3),
                     title, font_size=14, color=CREAM, bold=True)

        # Description
        add_text_box(slide, MARGIN + Inches(0.55), y + Inches(0.38),
                     card_w - Inches(0.8), Inches(0.45),
                     desc, font_size=9.5, color=LIGHT_GRAY, bold=False, line_spacing=13)

        # Source
        add_text_box(slide, MARGIN + Inches(0.55), y + Inches(0.8),
                     Inches(4), Inches(0.2),
                     f"Source: {src}", font_size=7, color=DARK_GRAY, bold=False)


def build_slide_13_close(prs):
    """Slide 13 — Close."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_gold_bar(slide)

    # ALCHE centered, large
    add_text_box(slide, Inches(0), Inches(1.6), SLIDE_W, Inches(1.0),
                 "ALCHE", font_size=64, color=GOLD, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(0), Inches(2.7), SLIDE_W, Inches(0.5),
                 "The unoccupied position in longevity",
                 font_size=18, color=CREAM, bold=False,
                 alignment=PP_ALIGN.CENTER)

    # Divider
    line_w = Inches(2)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(4), Inches(3.4), line_w, Emu(9144))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    line.line.fill.background()

    # Date
    add_text_box(slide, Inches(0), Inches(3.7), SLIDE_W, Inches(0.3),
                 "Competitive Positioning Analysis  /  February 2026",
                 font_size=10, color=MUTED, bold=False,
                 alignment=PP_ALIGN.CENTER)


# ── Main build pipeline ───────────────────────────────────────────────────────

def build_deck():
    """Orchestrate full deck build."""
    base_dir = os.path.dirname(os.path.abspath(__file__))
    chart_path = os.path.join(base_dir, "positioning-map.png")
    output_path = os.path.join(base_dir, "alche-imperfect-longevity-positioning.pptx")

    print("Alche Positioning Deck Builder")
    print("=" * 40)

    # Step 1: Brand data
    print("\n1. Loading brand data...")
    brands = define_brand_data()
    print(f"   {len(brands)} brands loaded ({len(brands) - 1} competitors + ALCHE)")

    # Step 2: Generate chart
    print("\n2. Generating positioning map...")
    generate_positioning_map(brands, chart_path)

    # Step 3: Build slides
    print("\n3. Building slides...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        ("Slide  1: Title",                    lambda: build_slide_01_title(prs)),
        ("Slide  2: Retention Problem",        lambda: build_slide_02_retention(prs)),
        ("Slide  3: The Backlash",             lambda: build_slide_03_backlash(prs)),
        ("Slide  4: The Gender Gap",           lambda: build_slide_04_gender_gap(prs)),
        ("Slide  5: Positioning Map",          lambda: build_slide_05_positioning_map(prs, chart_path)),
        ("Slide  6: Optimization Quadrant",    lambda: build_slide_06_optimization_quadrant(prs)),
        ("Slide  7: Lifestyle Quadrant",       lambda: build_slide_07_lifestyle_quadrant(prs)),
        ("Slide  8: Anti-Establishment",       lambda: build_slide_08_anti_establishment(prs)),
        ("Slide  9: Proposed Position",        lambda: build_slide_09_proposed_position(prs)),
        ("Slide 10: Integration Moat",         lambda: build_slide_10_integration_moat(prs)),
        ("Slide 11: Market Sizing",            lambda: build_slide_11_market_sizing(prs)),
        ("Slide 12: Berlin Context",           lambda: build_slide_12_berlin(prs)),
        ("Slide 13: Close",                    lambda: build_slide_13_close(prs)),
    ]

    for label, builder in builders:
        print(f"   {label}")
        builder()

    # Step 4: Save
    print(f"\n4. Saving deck...")
    prs.save(output_path)
    print(f"   Saved: {output_path}")

    print(f"\n{'=' * 40}")
    print(f"Done. {len(prs.slides)} slides generated.")
    print(f"  Deck:  {output_path}")
    print(f"  Chart: {chart_path}")


if __name__ == "__main__":
    build_deck()
