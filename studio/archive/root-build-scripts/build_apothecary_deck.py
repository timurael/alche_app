#!/usr/bin/env python3
"""
Alche "Imperfect Longevity" — The Apothecary Table Theme
Aesop meets Nature Medicine. Glass-morphism cards, warm amber on dark charcoal.
"""

import os
import glob

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import matplotlib.patheffects as pe
import numpy as np
from PIL import Image
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Font setup ─────────────────────────────────────────────────────────────────

FONT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fonts")

def setup_fonts():
    """Register downloaded fonts with matplotlib."""
    registered = {"serif": None, "sans": None}
    for ttf in glob.glob(os.path.join(FONT_DIR, "*.ttf")):
        try:
            fm.fontManager.addfont(ttf)
            fobj = fm.get_font(ttf)
            name = fobj.family_name
            if "Cormorant" in name and registered["serif"] is None:
                registered["serif"] = name
            elif "Outfit" in name and registered["sans"] is None:
                registered["sans"] = name
        except Exception:
            pass

    SERIF = registered["serif"] or "Baskerville"
    SANS = registered["sans"] or "Avenir"
    print(f"  Fonts: serif={SERIF}, sans={SANS}")
    return SERIF, SANS

MPL_SERIF, MPL_SANS = setup_fonts()

# PPTX font names (will reference installed fonts)
FONT_SERIF = "Cormorant Garamond"
FONT_SANS = "Outfit"
# Fallbacks if not installed on viewing machine:
FONT_SERIF_FB = "Baskerville"
FONT_SANS_FB = "Avenir Next"

# ── Design tokens ──────────────────────────────────────────────────────────────

BG          = RGBColor(0x2C, 0x24, 0x18)
BG_HEX      = "#2C2418"
GLASS_FILL  = RGBColor(0x40, 0x36, 0x2C)
GLASS_BORDER = RGBColor(0x5E, 0x52, 0x43)
AMBER       = RGBColor(0xD4, 0xA8, 0x53)
AMBER_HEX   = "#D4A853"
AMBER_LIGHT = RGBColor(0xE8, 0xC9, 0x7A)
AMBER_MUTED = RGBColor(0x8B, 0x74, 0x42)
CREAM       = RGBColor(0xF5, 0xED, 0xE0)
CREAM_HEX   = "#F5EDE0"
STONE       = RGBColor(0xB0, 0xA4, 0x8F)
STONE_HEX   = "#B0A48F"
STONE_DARK  = RGBColor(0x7A, 0x6F, 0x5E)
WARM_GRAY   = RGBColor(0x4A, 0x3F, 0x33)
WARM_DARK   = RGBColor(0x38, 0x2F, 0x24)

SLIDE_W     = Inches(10)
SLIDE_H     = Inches(5.625)
MARGIN      = Inches(0.7)
AMBER_BAR_H = Emu(18288)  # 2px amber rule at top

# Letter spacing for labels (hundredths of a point)
LABEL_SPACING = 250
SECTION_SPACING = 350

# ── Brand data ─────────────────────────────────────────────────────────────────

def define_brand_data():
    return [
        {"name": "Blueprint",       "x": 9.5, "y": 1.0, "size": 1.5, "color": 2},
        {"name": "Function Health", "x": 8.5, "y": 3.0, "size": 3.0, "color": 3},
        {"name": "Ultrahuman",      "x": 8.0, "y": 2.5, "size": 2.0, "color": 3},
        {"name": "Fountain Life",   "x": 8.5, "y": 2.0, "size": 2.5, "color": 2},
        {"name": "NOVOS Labs",      "x": 7.5, "y": 3.5, "size": 3.0, "color": 2},
        {"name": "Superpower",      "x": 7.0, "y": 3.0, "size": 2.5, "color": 2},
        {"name": "WHOOP",           "x": 7.5, "y": 2.5, "size": 2.0, "color": 3},
        {"name": "Oura",            "x": 7.0, "y": 4.0, "size": 4.5, "color": 3},
        {"name": "Eight Sleep",     "x": 7.0, "y": 3.5, "size": 3.0, "color": 3},
        {"name": "Levels",          "x": 7.5, "y": 3.0, "size": 2.5, "color": 3},
        {"name": "Apple Health",    "x": 5.0, "y": 6.0, "size": 7.0, "color": 1},
        {"name": "Parsley Health",  "x": 7.0, "y": 5.0, "size": 6.0, "color": 1},
        {"name": "Equinox",         "x": 4.0, "y": 4.5, "size": 4.0, "color": 3},
        {"name": "Life Time",       "x": 4.0, "y": 5.5, "size": 5.5, "color": 4},
        {"name": "Noom",            "x": 5.5, "y": 5.5, "size": 6.5, "color": 3},
        {"name": "Goop",            "x": 2.5, "y": 7.0, "size": 7.0, "color": 3},
        {"name": "Erewhon",         "x": 1.5, "y": 7.5, "size": 5.0, "color": 3},
        {"name": "Aman Wellness",   "x": 3.0, "y": 7.0, "size": 4.0, "color": 3},
        {"name": "Six Senses",      "x": 3.0, "y": 6.5, "size": 4.0, "color": 2},
        {"name": "AG1",             "x": 4.0, "y": 5.0, "size": 5.0, "color": 3},
        {"name": "Seed",            "x": 6.0, "y": 4.0, "size": 4.5, "color": 2},
        {"name": "Thorne",          "x": 6.5, "y": 2.5, "size": 3.0, "color": 1},
        {"name": "Moon Juice",      "x": 2.0, "y": 6.5, "size": 6.0, "color": 2},
        {"name": "Momentous",       "x": 6.0, "y": 3.0, "size": 2.5, "color": 1},
        {"name": "Liquid Death",    "x": 0.5, "y": 9.0, "size": 7.5, "color": 1},
        {"name": "MUD\\WTR",        "x": 2.0, "y": 8.5, "size": 6.5, "color": 1},
        {"name": "ANTI Berlin",     "x": 3.0, "y": 8.0, "size": 5.5, "color": 2},
        {"name": "Daluma Berlin",   "x": 2.5, "y": 7.0, "size": 5.0, "color": 2},
        {"name": "OneSkin",         "x": 7.0, "y": 5.0, "size": 6.5, "color": 1},
        {"name": "Timeline",        "x": 7.5, "y": 3.5, "size": 4.0, "color": 1},
        {"name": "ALCHE",           "x": 8.0, "y": 8.5, "size": 9.0, "color": 6},
    ]


# ── Chart generation ───────────────────────────────────────────────────────────

def score_to_color(score):
    """Integration score 1-6 → warm color gradient."""
    t = (score - 1) / 5.0
    r = int(107 + t * (212 - 107))
    g = int(80 + t * (168 - 80))
    b = int(64 + t * (83 - 64))
    return (r / 255, g / 255, b / 255)


def generate_positioning_map(brands, path):
    """Create the 4-axis positioning scatter chart with apothecary theme."""
    fig, ax = plt.subplots(figsize=(12, 7.5))
    fig.patch.set_facecolor(BG_HEX)
    ax.set_facecolor(BG_HEX)

    # Quadrant dividers
    ax.axhline(y=5.0, color="#4A3F33", linewidth=0.6, linestyle="-", zorder=1)
    ax.axvline(x=5.0, color="#4A3F33", linewidth=0.6, linestyle="-", zorder=1)

    # Dashed amber rectangle for unoccupied quadrant
    from matplotlib.patches import Rectangle
    rect = Rectangle((5.0, 5.0), 5.5, 5.5, linewidth=1.5, edgecolor=AMBER_HEX,
                      facecolor=AMBER_HEX, alpha=0.04, linestyle="--", zorder=1)
    ax.add_patch(rect)

    # Quadrant labels
    lp = dict(fontsize=8.5, fontfamily=MPL_SANS, alpha=0.30, ha="center")
    ax.text(2.5, 0.4, "UNDIFFERENTIATED", color=STONE_HEX, **lp)
    ax.text(7.5, 0.4, "OPTIMIZATION-FIRST", color=STONE_HEX, **lp)
    ax.text(2.5, 9.8, "LIFESTYLE WELLNESS", color=STONE_HEX, **lp)
    ax.text(7.5, 9.8, "UNOCCUPIED", color=AMBER_HEX, fontsize=10,
            fontfamily=MPL_SERIF, alpha=0.6, ha="center", fontstyle="italic")

    # Plot brands
    for b in brands:
        if b["name"] == "ALCHE":
            continue
        s = (b["size"] * 18) ** 1.6
        c = score_to_color(b["color"])
        ax.scatter(b["x"], b["y"], s=s, c=[c], alpha=0.6, edgecolors="none", zorder=3)

        # Label offsets
        ox, oy = 0.18, -0.28
        name = b["name"]
        if name in ("NOVOS Labs", "Timeline"):
            oy = 0.28
        elif name == "Levels":
            ox, oy = 0.25, 0.2
        elif name == "Eight Sleep":
            oy = -0.35
        elif name == "Superpower":
            oy = 0.25
        elif name == "Daluma Berlin":
            ox, oy = -0.3, -0.5
        elif name == "OneSkin":
            ox, oy = -0.2, 0.3
        elif name == "Parsley Health":
            ox, oy = -0.3, -0.3
        elif name == "Aman Wellness":
            ox, oy = 0.3, -0.45
        elif name == "Goop":
            ox, oy = -0.4, -0.05
        elif name == "Six Senses":
            ox, oy = 0.3, -0.1

        ax.text(b["x"] + ox, b["y"] + oy, name,
                fontsize=6.5, color=STONE_HEX, fontfamily=MPL_SANS, alpha=0.75, zorder=4)

    # ALCHE with glow
    alche = next(b for b in brands if b["name"] == "ALCHE")
    ax.scatter(alche["x"], alche["y"], s=2800, c=[AMBER_HEX], alpha=0.10, edgecolors="none", zorder=5)
    ax.scatter(alche["x"], alche["y"], s=1600, c=[AMBER_HEX], alpha=0.18, edgecolors="none", zorder=5)
    ax.scatter(alche["x"], alche["y"], s=900, c=[AMBER_HEX], alpha=0.85,
               edgecolors=CREAM_HEX, linewidth=1.5, zorder=6)
    ax.text(alche["x"] + 0.35, alche["y"] + 0.05, "ALCHE", fontsize=12, color=CREAM_HEX,
            fontfamily=MPL_SERIF, fontweight="bold", zorder=7,
            path_effects=[pe.withStroke(linewidth=3, foreground=BG_HEX)])

    # Axes styling
    ax.set_xlim(-0.3, 10.5)
    ax.set_ylim(-0.3, 10.5)
    ax.set_xlabel("Lifestyle / Aesthetic  \u2190\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2192  Clinical / Data-Driven",
                  fontsize=8.5, color="#7A6F5E", fontfamily=MPL_SANS, labelpad=10)
    ax.set_ylabel("Ascetic / Optimization  \u2190\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2014\u2192  Human / Joyful",
                  fontsize=8.5, color="#7A6F5E", fontfamily=MPL_SANS, labelpad=10)

    ax.tick_params(colors="#5E5243", labelsize=7)
    for spine in ax.spines.values():
        spine.set_color("#4A3F33")
        spine.set_linewidth(0.5)
    ax.set_xticks(range(0, 11))
    ax.set_yticks(range(0, 11))
    ax.grid(True, alpha=0.06, color="#5E5243", linewidth=0.4)

    # Size legend
    for val, label in [(3, "Low inclusion"), (6, "Medium"), (9, "High inclusion")]:
        s = (val * 18) ** 1.6
        ax.scatter([], [], s=s, c=["#8B7442"], alpha=0.45, edgecolors="none", label=label)
    leg = ax.legend(loc="lower left", frameon=False, fontsize=7,
                    labelcolor=STONE_HEX, handletextpad=1.2, borderpad=1.0,
                    title="Bubble size = Inclusivity", title_fontproperties={"size": 7.5, "family": MPL_SANS})
    leg.get_title().set_color(STONE_HEX)

    # Color bar
    from matplotlib.colors import LinearSegmentedColormap
    cmap_colors = [score_to_color(i) for i in range(1, 7)]
    cmap = LinearSegmentedColormap.from_list("integration", cmap_colors, N=6)
    sm = plt.cm.ScalarMappable(cmap=cmap, norm=plt.Normalize(vmin=1, vmax=6))
    sm.set_array([])
    cbar_ax = fig.add_axes([0.78, 0.08, 0.15, 0.015])
    cbar = fig.colorbar(sm, cax=cbar_ax, orientation="horizontal", ticks=[1, 2, 3, 4, 5, 6])
    cbar.ax.tick_params(labelsize=6, colors=STONE_HEX, length=2)
    cbar.outline.set_edgecolor("#4A3F33")
    cbar.outline.set_linewidth(0.5)
    cbar.set_label("Integration pillars (1\u20136)", fontsize=7, color=STONE_HEX,
                   fontfamily=MPL_SANS, labelpad=4)

    fig.subplots_adjust(left=0.08, right=0.95, top=0.95, bottom=0.12)
    fig.savefig(path, dpi=300, facecolor=BG_HEX, bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)
    print(f"  Chart saved: {path}")


# ── Slide helpers ──────────────────────────────────────────────────────────────

def set_slide_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_amber_rule(slide):
    """Thin 2px amber rule at top of slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, AMBER_BAR_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = AMBER
    shape.line.fill.background()

def _set_fill_transparency(shape, transparency_pct):
    """Set fill transparency. 0=opaque, 100=invisible."""
    opacity_val = str(int((100 - transparency_pct) * 1000))
    spPr = shape._element.spPr
    solid_fill = spPr.find(qn('a:solidFill'))
    if solid_fill is not None:
        color_elem = solid_fill[0]
        for a in list(color_elem):
            if a.tag == qn('a:alpha'):
                color_elem.remove(a)
        alpha = etree.SubElement(color_elem, qn('a:alpha'))
        alpha.set('val', opacity_val)

def _set_letter_spacing(run, spacing_hundredths):
    """Set character spacing in hundredths of a point."""
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(spacing_hundredths))

def add_glass_panel(slide, left, top, width, height, border_color=GLASS_BORDER,
                    transparency=30, corner_radius=Inches(0.08)):
    """Glass-morphism card: translucent fill with subtle border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GLASS_FILL
    _set_fill_transparency(shape, transparency)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_glass_vial(slide, left, top, width, height, stat_text, label_text,
                   stat_size=72, stat_color=AMBER, label_color=STONE,
                   source_text=None):
    """Glass vial motif: translucent container with amber cap and data inside."""
    # Vial body
    panel = add_glass_panel(slide, left, top, width, height, border_color=AMBER_MUTED)

    # Amber cap — thin line at top inside the vial
    cap_h = Emu(27432)  # 3px
    cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 left + Inches(0.06), top + Inches(0.06),
                                 width - Inches(0.12), cap_h)
    cap.fill.solid()
    cap.fill.fore_color.rgb = AMBER
    cap.line.fill.background()
    _set_fill_transparency(cap, 30)

    # Stat number
    add_text(slide, left + Inches(0.15), top + Inches(0.2),
             width - Inches(0.3), Inches(0.8),
             stat_text, font_size=stat_size, color=stat_color,
             bold=False, font_name=FONT_SERIF, alignment=PP_ALIGN.CENTER)

    # Label
    add_text(slide, left + Inches(0.1), top + height - Inches(0.7),
             width - Inches(0.2), Inches(0.5),
             label_text, font_size=9.5, color=label_color,
             bold=False, font_name=FONT_SANS, alignment=PP_ALIGN.CENTER,
             line_spacing=13)

    # Source
    if source_text:
        add_text(slide, left + Inches(0.1), top + height - Inches(0.25),
                 width - Inches(0.2), Inches(0.2),
                 source_text, font_size=6.5, color=STONE_DARK,
                 bold=False, font_name=FONT_SANS, alignment=PP_ALIGN.CENTER)

    return panel


def add_text(slide, left, top, width, height, text, font_size=14,
             color=CREAM, bold=False, font_name=FONT_SANS,
             alignment=PP_ALIGN.LEFT, line_spacing=None, letter_spacing=None):
    """Add a text box. Returns the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.bold = bold
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    if letter_spacing:
        for run in p.runs:
            _set_letter_spacing(run, letter_spacing)
    return txBox


def add_section_label(slide, number, label):
    """Section label: '01 — LABEL' in spaced Outfit."""
    text = f"{number:02d}  \u2014  {label}"
    txBox = add_text(slide, MARGIN, Inches(0.3), Inches(5), Inches(0.22),
                     text, font_size=8.5, color=STONE_DARK,
                     font_name=FONT_SANS, letter_spacing=SECTION_SPACING)
    return txBox


def add_section_number(slide, number):
    """Faint section number in bottom-right."""
    add_text(slide, Inches(8.5), Inches(4.85), Inches(0.8), Inches(0.4),
             f"{number:02d}", font_size=22, color=WARM_GRAY,
             font_name=FONT_SERIF, alignment=PP_ALIGN.RIGHT)


def add_thin_rule(slide, left, top, width):
    """Thin amber-tinted divider line."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(6858))
    line.fill.solid()
    line.fill.fore_color.rgb = WARM_GRAY
    line.line.fill.background()
    return line


def add_headline(slide, text, top=Inches(0.58), width=Inches(8.5),
                 font_size=44, color=CREAM, bold=False):
    """Oversized Cormorant Garamond headline."""
    return add_text(slide, MARGIN, top, width, Inches(1.0),
                    text, font_size=font_size, color=color,
                    bold=bold, font_name=FONT_SERIF, line_spacing=font_size + 6)


def add_multi_run(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT,
                  line_spacing=None):
    """Text box with multiple styled runs: [(text, size, color, bold, font_name), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    for i, (text, size, color, bold, fname) in enumerate(runs):
        run = p.runs[0] if i == 0 and p.runs else p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = fname
        run.font.bold = bold
    return txBox


def add_source(slide, text, left=None, top=Inches(5.1)):
    """Source citation at bottom."""
    if left is None:
        left = MARGIN
    return add_text(slide, left, top, Inches(7), Inches(0.2),
                    text, font_size=6.5, color=STONE_DARK, font_name=FONT_SANS)


# ── Slide builders ─────────────────────────────────────────────────────────────

def build_slide_01_title(prs):
    """Title — hero text with apothecary warmth."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)

    # ALCHE wordmark — spaced
    txBox = add_text(slide, MARGIN, Inches(0.45), Inches(3), Inches(0.35),
                     "ALCHE", font_size=14, color=AMBER,
                     font_name=FONT_SANS, bold=False, letter_spacing=500)

    # Headline
    add_multi_run(slide, MARGIN, Inches(1.0), Inches(8.2), Inches(1.8), [
        ("No longevity brand occupies\nthe intersection of deep science ", 52, CREAM, False, FONT_SERIF),
        ("and human joy", 52, AMBER, False, FONT_SERIF),
    ], line_spacing=60)

    # Subtitle
    add_text(slide, MARGIN, Inches(2.85), Inches(7), Inches(0.35),
             "Evidence from 30 brands across the $6.8T wellness economy",
             font_size=12, color=STONE, font_name=FONT_SANS)

    # Divider
    add_thin_rule(slide, MARGIN, Inches(3.4), Inches(8.0))

    # Three glass vials in stat row
    vial_w = Inches(2.35)
    vial_h = Inches(1.7)
    vial_gap = Inches(0.25)
    vial_y = Inches(3.6)

    stats = [("30", "brands\nanalyzed"), ("4", "positioning\naxes"), ("0", "in the\nwhite space")]
    for i, (val, label) in enumerate(stats):
        x = MARGIN + i * (vial_w + vial_gap)
        add_glass_vial(slide, x, vial_y, vial_w, vial_h,
                       val, label, stat_size=52, stat_color=AMBER)

    # Bottom credit
    add_text(slide, MARGIN, Inches(5.1), Inches(5), Inches(0.2),
             "Competitive Positioning Analysis  \u00b7  February 2026",
             font_size=7.5, color=STONE_DARK, font_name=FONT_SANS,
             letter_spacing=150)


def build_slide_02_retention(prs):
    """The Retention Problem — large stat callout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)
    add_section_label(slide, 1, "THE RETENTION PROBLEM")
    add_section_number(slide, 1)

    add_headline(slide, "90% of wearable users abandon their\ndevice within 12 months",
                 font_size=38, top=Inches(0.6))

    # Large vial with 90%
    add_glass_vial(slide, MARGIN, Inches(1.7), Inches(3.8), Inches(2.8),
                   "90%", "wearable\nabandonment rate",
                   stat_size=96, stat_color=AMBER)

    # Smaller vial for 3.4%
    add_glass_vial(slide, Inches(5.2), Inches(1.7), Inches(2.8), Inches(1.7),
                   "3.4%", "health app retention\nat day 30",
                   stat_size=56, stat_color=AMBER_LIGHT)

    # Takeaway in glass panel
    panel = add_glass_panel(slide, Inches(5.2), Inches(3.65), Inches(4.1), Inches(0.85))
    add_text(slide, Inches(5.4), Inches(3.75), Inches(3.7), Inches(0.65),
             "The current model optimizes for measurement.\nUsers leave anyway.",
             font_size=12.5, color=CREAM, font_name=FONT_SERIF,
             line_spacing=18)

    add_source(slide, "Sources: Industry research data on wearable retention and health app engagement")


def build_slide_03_backlash(prs):
    """The Backlash — quote + evidence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)
    add_section_label(slide, 2, "THE BACKLASH")
    add_section_number(slide, 2)

    add_headline(slide, "The Global Wellness Summit named the\nover-optimization backlash its #2 trend for 2026",
                 font_size=34)

    # Left: quote in glass panel (wide vial)
    panel = add_glass_panel(slide, MARGIN, Inches(1.65), Inches(4.5), Inches(2.5),
                            border_color=AMBER_MUTED)
    # Amber cap
    cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 MARGIN + Inches(0.06), Inches(1.71),
                                 Inches(4.38), Emu(27432))
    cap.fill.solid()
    cap.fill.fore_color.rgb = AMBER
    cap.line.fill.background()
    _set_fill_transparency(cap, 30)

    # Quote mark
    add_text(slide, MARGIN + Inches(0.2), Inches(1.85), Inches(0.5), Inches(0.5),
             "\u201C", font_size=42, color=AMBER, font_name=FONT_SERIF)
    add_text(slide, MARGIN + Inches(0.2), Inches(2.2), Inches(4.0), Inches(1.2),
             "Wellness experiences will embrace what humans actually are: imperfect, emotional, relational and sensory.",
             font_size=14.5, color=CREAM, font_name=FONT_SERIF,
             line_spacing=21)
    add_text(slide, MARGIN + Inches(0.2), Inches(3.5), Inches(4.0), Inches(0.3),
             "\u2014 Global Wellness Summit, Future of Wellness 2026",
             font_size=8, color=STONE_DARK, font_name=FONT_SANS,
             letter_spacing=100)

    # Right: three evidence points
    rx = Inches(5.6)
    points = [
        ("Trend #2 of 2026", "Rehumanizing wellness \u2014 from metrics-obsessed to experience-centered"),
        ("Consumer signal", "Data-overload fatigue driving users away is cultural, not technical"),
        ("Market implication", "Brands combining scientific rigor with emotional resonance will capture the backlash"),
    ]
    for i, (head, body) in enumerate(points):
        y = Inches(1.7) + Inches(i * 0.95)
        add_text(slide, rx, y, Inches(3.7), Inches(0.22),
                 head, font_size=10, color=AMBER, font_name=FONT_SANS,
                 bold=False, letter_spacing=LABEL_SPACING)
        add_text(slide, rx, y + Inches(0.25), Inches(3.7), Inches(0.55),
                 body, font_size=9.5, color=STONE, font_name=FONT_SANS,
                 line_spacing=14)

    add_source(slide, "Sources: GWS Future of Wellness 2026, PR Newswire, Global Wellness Institute")


def build_slide_04_gender_gap(prs):
    """The Gender Gap — three stat vials + context."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)
    add_section_label(slide, 3, "THE GENDER GAP")
    add_section_number(slide, 3)

    add_headline(slide, "Women\u2019s health receives 6% of private\nhealthcare investment despite being 51%\nof the population",
                 font_size=32)

    # Three stat vials
    vial_w = Inches(2.5)
    vial_h = Inches(2.05)
    gap = Inches(0.15)
    vy = Inches(1.95)

    stats = [
        ("6%", "of private healthcare\ninvestment goes to\nwomen\u2019s health", "WEF"),
        ("$40B", "menopause market\nby 2030 (8\u00d7 growth)", "BCG"),
        ("$100M", "Melinda French Gates\ncommitment to\nwomen\u2019s health R&D", "Pivotal Ventures"),
    ]
    for i, (val, label, src) in enumerate(stats):
        x = MARGIN + i * (vial_w + gap)
        add_glass_vial(slide, x, vy, vial_w, vial_h,
                       val, label, stat_size=42, source_text=f"Source: {src}")

    # Quote panel on right
    qx = Inches(5.6)
    qy = Inches(4.2)
    panel = add_glass_panel(slide, MARGIN, qy, Inches(8.6), Inches(0.7), border_color=AMBER_MUTED)
    add_text(slide, MARGIN + Inches(0.15), qy + Inches(0.08), Inches(0.3), Inches(0.3),
             "\u201C", font_size=28, color=AMBER, font_name=FONT_SERIF)
    add_text(slide, MARGIN + Inches(0.35), qy + Inches(0.12), Inches(7.5), Inches(0.5),
             "The booming longevity market, like medicine decades before it, is tacitly male.  \u2014 GWS 2026",
             font_size=12, color=CREAM, font_name=FONT_SERIF, line_spacing=17)


def build_slide_05_map(prs, chart_path):
    """THE POSITIONING MAP — centerpiece."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)
    add_section_number(slide, 4)

    add_section_label(slide, 4, "POSITIONING MAP")
    add_headline(slide, "30 brands plotted on 4 axes \u2014 the top-right quadrant is empty",
                 font_size=24, top=Inches(0.45))

    # Embed chart
    img = Image.open(chart_path)
    aspect = img.height / img.width
    max_h = 4.5
    cw = min(9.6, max_h / aspect)
    ch = cw * aspect
    cx = (10 - cw) / 2
    slide.shapes.add_picture(chart_path, Inches(cx), Inches(0.92), Inches(cw), Inches(ch))


def build_slide_06_optimization(prs):
    """Optimization-First Quadrant — brand cards in glass panels."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)
    add_section_label(slide, 5, "OPTIMIZATION-FIRST QUADRANT")
    add_section_number(slide, 5)

    add_headline(slide, "High-science brands default to ascetic\npositioning and male-oriented design",
                 font_size=32)

    cards = [
        ("Blueprint", "Science 9.5  \u00b7  Humanity 1.0  \u00b7  Inclusion 1.5",
         "Bryan Johnson\u2019s protocol-driven approach. Highest science, lowest humanity in dataset."),
        ("Function Health", "Science 8.5  \u00b7  Humanity 3.0  \u00b7  Inclusion 3.0",
         "100+ biomarker testing platform. Clinical positioning, limited lifestyle integration."),
        ("WHOOP", "Science 7.5  \u00b7  Humanity 2.5  \u00b7  Inclusion 2.0",
         "Performance-first wearable. Core audience: male athletes and biohackers."),
        ("Ultrahuman", "Science 8.0  \u00b7  Humanity 2.5  \u00b7  Inclusion 2.0",
         "Metabolic tracking via CGM. Tech-forward language, narrow demographic appeal."),
    ]

    cw = Inches(4.05)
    ch = Inches(1.45)
    gx, gy_gap = Inches(0.25), Inches(0.2)
    sy = Inches(1.6)

    for i, (name, scores, desc) in enumerate(cards):
        col, row = i % 2, i // 2
        x = MARGIN + col * (cw + gx)
        y = sy + row * (ch + gy_gap)
        add_glass_panel(slide, x, y, cw, ch)

        add_text(slide, x + Inches(0.15), y + Inches(0.1), cw - Inches(0.3), Inches(0.28),
                 name, font_size=15, color=CREAM, font_name=FONT_SERIF, bold=False)
        add_text(slide, x + Inches(0.15), y + Inches(0.38), cw - Inches(0.3), Inches(0.22),
                 scores, font_size=9, color=AMBER, font_name=FONT_SANS,
                 letter_spacing=80)
        add_text(slide, x + Inches(0.15), y + Inches(0.65), cw - Inches(0.3), Inches(0.65),
                 desc, font_size=9, color=STONE, font_name=FONT_SANS, line_spacing=13)


def build_slide_07_lifestyle(prs):
    """Lifestyle Wellness Quadrant — brand cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)
    add_section_label(slide, 6, "LIFESTYLE WELLNESS QUADRANT")
    add_section_number(slide, 6)

    add_headline(slide, "High-humanity brands lack clinical depth\nor longevity specificity",
                 font_size=32)

    cards = [
        ("Goop", "Science 2.5  \u00b7  Humanity 7.0  \u00b7  Inclusion 7.0",
         "Lifestyle-first wellness brand. Strong cultural resonance, limited scientific credibility."),
        ("Erewhon", "Science 1.5  \u00b7  Humanity 7.5  \u00b7  Inclusion 5.0",
         "Premium grocery as lifestyle brand. High aspiration, minimal health-outcome depth."),
        ("Moon Juice", "Science 2.0  \u00b7  Humanity 6.5  \u00b7  Inclusion 6.0",
         "Adaptogen-forward product brand. Aesthetic strength without clinical grounding."),
        ("Aman Wellness", "Science 3.0  \u00b7  Humanity 7.0  \u00b7  Inclusion 4.0",
         "Ultra-luxury wellness hospitality. Beautiful experience, low science integration."),
    ]

    cw = Inches(4.05)
    ch = Inches(1.45)
    gx, gy_gap = Inches(0.25), Inches(0.2)
    sy = Inches(1.6)

    for i, (name, scores, desc) in enumerate(cards):
        col, row = i % 2, i // 2
        x = MARGIN + col * (cw + gx)
        y = sy + row * (ch + gy_gap)
        add_glass_panel(slide, x, y, cw, ch)

        add_text(slide, x + Inches(0.15), y + Inches(0.1), cw - Inches(0.3), Inches(0.28),
                 name, font_size=15, color=CREAM, font_name=FONT_SERIF)
        add_text(slide, x + Inches(0.15), y + Inches(0.38), cw - Inches(0.3), Inches(0.22),
                 scores, font_size=9, color=AMBER, font_name=FONT_SANS,
                 letter_spacing=80)
        add_text(slide, x + Inches(0.15), y + Inches(0.65), cw - Inches(0.3), Inches(0.65),
                 desc, font_size=9, color=STONE, font_name=FONT_SANS, line_spacing=13)


def build_slide_08_anti_establishment(prs):
    """Anti-Establishment Precedent — two glass case studies."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)
    add_section_label(slide, 7, "CATEGORY INVERSION PRECEDENT")
    add_section_number(slide, 7)

    add_headline(slide, "Liquid Death reached $1.4B valuation by\ninverting category conventions in water",
                 font_size=32)

    case_w = Inches(4.0)
    case_h = Inches(2.7)

    # Liquid Death panel
    lx = MARGIN
    add_glass_panel(slide, lx, Inches(1.65), case_w, case_h, border_color=AMBER_MUTED)
    # Amber cap
    cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 lx + Inches(0.06), Inches(1.71), case_w - Inches(0.12), Emu(27432))
    cap.fill.solid()
    cap.fill.fore_color.rgb = AMBER
    cap.line.fill.background()
    _set_fill_transparency(cap, 30)

    add_text(slide, lx + Inches(0.2), Inches(1.85), Inches(3.5), Inches(0.3),
             "Liquid Death", font_size=20, color=CREAM, font_name=FONT_SERIF)
    add_text(slide, lx + Inches(0.2), Inches(2.18), Inches(3.5), Inches(0.22),
             "Commoditized water \u2192 punk aesthetic", font_size=9, color=AMBER,
             font_name=FONT_SANS, letter_spacing=80)

    for j, (val, lab) in enumerate([("$1.4B", "valuation"), ("$333M", "revenue (2024)")]):
        sy = Inches(2.55) + Inches(j * 0.6)
        add_text(slide, lx + Inches(0.2), sy, Inches(1.5), Inches(0.4),
                 val, font_size=30, color=AMBER, font_name=FONT_SERIF)
        add_text(slide, lx + Inches(1.8), sy + Inches(0.08), Inches(2), Inches(0.3),
                 lab, font_size=9.5, color=STONE, font_name=FONT_SANS)

    add_text(slide, lx + Inches(0.2), Inches(3.85), Inches(3.5), Inches(0.2),
             "Sources: Fueler, Sacra, Entrepreneur", font_size=6.5, color=STONE_DARK,
             font_name=FONT_SANS)

    # MUD\WTR panel
    rx = Inches(5.2)
    add_glass_panel(slide, rx, Inches(1.65), case_w, case_h, border_color=AMBER_MUTED)
    cap2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  rx + Inches(0.06), Inches(1.71), case_w - Inches(0.12), Emu(27432))
    cap2.fill.solid()
    cap2.fill.fore_color.rgb = AMBER
    cap2.line.fill.background()
    _set_fill_transparency(cap2, 30)

    add_text(slide, rx + Inches(0.2), Inches(1.85), Inches(3.5), Inches(0.3),
             "MUD\\WTR", font_size=20, color=CREAM, font_name=FONT_SERIF)
    add_text(slide, rx + Inches(0.2), Inches(2.18), Inches(3.5), Inches(0.22),
             "Anti-hustle-culture coffee alternative", font_size=9, color=AMBER,
             font_name=FONT_SANS, letter_spacing=80)

    add_text(slide, rx + Inches(0.2), Inches(2.55), Inches(1.5), Inches(0.4),
             "$47M", font_size=30, color=AMBER, font_name=FONT_SERIF)
    add_text(slide, rx + Inches(1.8), Inches(2.63), Inches(2), Inches(0.3),
             "revenue (2024 est.)", font_size=9.5, color=STONE, font_name=FONT_SANS)

    add_text(slide, rx + Inches(0.2), Inches(3.85), Inches(3.5), Inches(0.2),
             "Sources: ecdb, Owler", font_size=6.5, color=STONE_DARK, font_name=FONT_SANS)

    # Bottom takeaway
    add_text(slide, MARGIN, Inches(4.6), Inches(8.2), Inches(0.5),
             "Neither operates in longevity. The archetype is proven but the category is untouched.",
             font_size=13, color=CREAM, font_name=FONT_SERIF)


def build_slide_09_proposed_position(prs):
    """The Proposed Position — three glass vial pillars."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)
    add_section_label(slide, 8, "THE PROPOSED POSITION")
    add_section_number(slide, 8)

    add_headline(slide, "Alche sits at the intersection: clinical-grade\nscience, inclusive design, joyful brand",
                 font_size=32)

    col_w = Inches(2.65)
    col_h = Inches(3.15)
    gap = Inches(0.2)
    sy = Inches(1.65)

    columns = [
        ("Clinical Science", "8.0 / 10", [
            "Bloodwork & biomarker panels",
            "AI-personalized protocols",
            "Evidence-based supplementation",
            "Longitudinal health tracking",
        ]),
        ("Human Brand", "8.5 / 10", [
            "Permission-giving, not punishing",
            "Lifestyle integration over restriction",
            "Joyful aesthetics & cultural edge",
            "Community over competition",
        ]),
        ("Inclusive Design", "9.0 / 10", [
            "Women-first approach",
            "Hormonal health integration",
            "Diverse bodies & life stages",
            "Accessible entry points",
        ]),
    ]

    for i, (title, score, bullets) in enumerate(columns):
        x = MARGIN + i * (col_w + gap)

        # Glass vial panel
        add_glass_panel(slide, x, sy, col_w, col_h, border_color=AMBER_MUTED)

        # Amber cap
        cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x + Inches(0.06), sy + Inches(0.06),
                                     col_w - Inches(0.12), Emu(27432))
        cap.fill.solid()
        cap.fill.fore_color.rgb = AMBER
        cap.line.fill.background()
        _set_fill_transparency(cap, 30)

        # Title
        add_text(slide, x + Inches(0.15), sy + Inches(0.15),
                 col_w - Inches(0.3), Inches(0.3),
                 title, font_size=16, color=AMBER, font_name=FONT_SERIF)

        # Score
        add_text(slide, x + Inches(0.15), sy + Inches(0.45),
                 col_w - Inches(0.3), Inches(0.22),
                 score, font_size=10, color=STONE, font_name=FONT_SANS,
                 letter_spacing=150)

        # Thin rule
        add_thin_rule(slide, x + Inches(0.15), sy + Inches(0.72), Inches(0.6))

        # Bullets
        for j, bullet in enumerate(bullets):
            add_text(slide, x + Inches(0.15), sy + Inches(0.9 + j * 0.48),
                     col_w - Inches(0.3), Inches(0.4),
                     f"\u2022  {bullet}", font_size=9.5, color=STONE,
                     font_name=FONT_SANS, line_spacing=13)


def build_slide_10_integration(prs):
    """Integration Moat — pillar comparison matrix."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)
    add_section_label(slide, 9, "INTEGRATION MOAT")
    add_section_number(slide, 9)

    add_headline(slide, "Alche scores 6/6 on integration pillars;\nthe highest competitor scores 4/6",
                 font_size=32)

    pillars = ["Physical\nSpace", "Products", "AI App", "Membership", "Community", "Integrated\nInfo System"]

    brands_compare = [
        ("ALCHE",     [1, 1, 1, 1, 1, 1], AMBER),
        ("Life Time", [1, 1, 1, 1, 0, 0], AMBER_MUTED),
        ("WHOOP",     [0, 1, 1, 1, 0, 0], RGBColor(0x6B, 0x5E, 0x48)),
        ("Equinox",   [1, 0, 1, 1, 0, 0], RGBColor(0x5E, 0x52, 0x40)),
    ]

    bar_start_x = Inches(2.2)
    bar_w_unit = Inches(1.05)
    bar_h = Inches(0.48)
    gap_y = Inches(0.12)
    sy = Inches(1.75)

    # Column headers
    for j, pillar in enumerate(pillars):
        px = bar_start_x + j * bar_w_unit + Inches(0.15)
        add_text(slide, px, Inches(1.42), bar_w_unit - Inches(0.1), Inches(0.35),
                 pillar, font_size=7, color=STONE_DARK, font_name=FONT_SANS,
                 alignment=PP_ALIGN.CENTER, letter_spacing=80)

    for i, (bname, scores, bcolor) in enumerate(brands_compare):
        y = sy + i * (bar_h + gap_y)
        is_alche = bname == "ALCHE"

        add_text(slide, MARGIN, y + Inches(0.06), Inches(1.4), Inches(0.35),
                 bname, font_size=13, color=CREAM if is_alche else STONE,
                 font_name=FONT_SERIF, bold=False)

        total = sum(scores)
        add_text(slide, Inches(8.6), y + Inches(0.06), Inches(0.7), Inches(0.35),
                 f"{total}/6", font_size=13, color=bcolor, font_name=FONT_SERIF,
                 alignment=PP_ALIGN.RIGHT)

        for j, has in enumerate(scores):
            cx = bar_start_x + j * bar_w_unit
            cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          cx, y, bar_w_unit - Inches(0.06), bar_h)
            if has:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bcolor
                if not is_alche:
                    _set_fill_transparency(cell, 20)
                cell.line.fill.background()
                add_text(slide, cx + Inches(0.3), y + Inches(0.06),
                         Inches(0.4), Inches(0.3),
                         "\u2713", font_size=15,
                         color=BG if is_alche else CREAM,
                         font_name=FONT_SANS, bold=True, alignment=PP_ALIGN.CENTER)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x34, 0x2C, 0x22)
                cell.line.fill.background()
                add_text(slide, cx + Inches(0.3), y + Inches(0.06),
                         Inches(0.4), Inches(0.3),
                         "\u2014", font_size=12, color=WARM_GRAY,
                         font_name=FONT_SANS, alignment=PP_ALIGN.CENTER)

    add_source(slide, "Source: Alche 6-pillar integration analysis of 30+ wellness & longevity companies",
               top=Inches(4.8))


def build_slide_11_market(prs):
    """Market Sizing — six glass vials in grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)
    add_section_label(slide, 10, "MARKET CONTEXT")
    add_section_number(slide, 10)

    add_headline(slide, "The wellness economy is $6.8T\nand growing to $10T by 2029",
                 font_size=34)

    stats = [
        ("$6.8T", "global wellness\neconomy (2025)", "Global Wellness Institute"),
        ("$10T", "projected\nby 2029", "Global Wellness Institute"),
        ("$40B", "menopause market\nby 2030", "BCG"),
        ("57%", "prioritize aging well\nmore than 5 yrs ago", "NielsenIQ"),
        ("84%", "say wellness is\na top priority", "McKinsey"),
        ("6%", "of private healthcare\ninvestment \u2192 women", "WEF"),
    ]

    vw = Inches(2.6)
    vh = Inches(1.55)
    gx, gy = Inches(0.17), Inches(0.15)
    sx = MARGIN
    sy = Inches(1.6)

    for i, (val, desc, src) in enumerate(stats):
        col, row = i % 3, i // 3
        x = sx + col * (vw + gx)
        y = sy + row * (vh + gy)
        add_glass_vial(slide, x, y, vw, vh,
                       val, desc, stat_size=34, source_text=f"Source: {src}")


def build_slide_12_berlin(prs):
    """Berlin Context — three horizontal glass panels."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)
    add_section_label(slide, 11, "BERLIN CONTEXT")
    add_section_number(slide, 11)

    add_headline(slide, "Berlin has an existing longevity scene and\ncultural DNA that aligns with the position",
                 font_size=32)

    points = [
        ("ANTI Berlin",
         "Longevity + culture hybrid on Brunnenstra\u00dfe. 20K Instagram followers. Single location proving demand for science-meets-lifestyle in Berlin.",
         "IGNANT, antispaces.com"),
        ("Berlin Biohackers",
         "2,213 members in the Berlin Biohackers meetup group. Active community of self-experimenters already engaged with longevity.",
         "Meetup / Reddit research"),
        ("LIFE Summit 2026",
         "Berlin hosting a dedicated longevity summit. The city is becoming a European hub for longevity discourse and innovation.",
         "LIFE Summit"),
    ]

    pw = Inches(8.6)
    ph = Inches(0.95)
    gap = Inches(0.12)
    sy = Inches(1.6)

    for i, (title, desc, src) in enumerate(points):
        y = sy + i * (ph + gap)
        add_glass_panel(slide, MARGIN, y, pw, ph, border_color=GLASS_BORDER)

        # Amber cap line inside
        cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     MARGIN + Inches(0.06), y + Inches(0.06),
                                     pw - Inches(0.12), Emu(18288))
        cap.fill.solid()
        cap.fill.fore_color.rgb = AMBER
        cap.line.fill.background()
        _set_fill_transparency(cap, 40)

        # Number
        add_text(slide, MARGIN + Inches(0.15), y + Inches(0.12), Inches(0.3), Inches(0.3),
                 str(i + 1), font_size=18, color=AMBER, font_name=FONT_SERIF,
                 alignment=PP_ALIGN.CENTER)

        # Title
        add_text(slide, MARGIN + Inches(0.55), y + Inches(0.1), Inches(3), Inches(0.25),
                 title, font_size=14, color=CREAM, font_name=FONT_SERIF)

        # Description
        add_text(slide, MARGIN + Inches(0.55), y + Inches(0.36), pw - Inches(0.8), Inches(0.38),
                 desc, font_size=9, color=STONE, font_name=FONT_SANS, line_spacing=13)

        # Source
        add_text(slide, MARGIN + Inches(0.55), y + Inches(0.73), Inches(3), Inches(0.15),
                 f"Source: {src}", font_size=6.5, color=STONE_DARK, font_name=FONT_SANS)


def build_slide_13_close(prs):
    """Close — centered, minimal, apothecary warmth."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_amber_rule(slide)

    # ALCHE centered large
    add_text(slide, Inches(0), Inches(1.4), SLIDE_W, Inches(1.0),
             "ALCHE", font_size=72, color=AMBER, font_name=FONT_SERIF,
             alignment=PP_ALIGN.CENTER, letter_spacing=600)

    # Subtitle
    add_text(slide, Inches(0), Inches(2.6), SLIDE_W, Inches(0.5),
             "The unoccupied position in longevity",
             font_size=18, color=CREAM, font_name=FONT_SERIF,
             alignment=PP_ALIGN.CENTER)

    # Thin amber rule
    rw = Inches(1.5)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(4.25), Inches(3.3), rw, Emu(9144))
    rule.fill.solid()
    rule.fill.fore_color.rgb = AMBER
    rule.line.fill.background()
    _set_fill_transparency(rule, 40)

    # Date
    add_text(slide, Inches(0), Inches(3.6), SLIDE_W, Inches(0.3),
             "Competitive Positioning Analysis  /  February 2026",
             font_size=9, color=STONE_DARK, font_name=FONT_SANS,
             alignment=PP_ALIGN.CENTER, letter_spacing=200)


# ── Main pipeline ──────────────────────────────────────────────────────────────

def build_deck():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    chart_path = os.path.join(base_dir, "positioning-map-apothecary.png")
    output_path = os.path.join(base_dir, "alche-apothecary-positioning.pptx")

    print("Alche Apothecary Table Deck Builder")
    print("=" * 45)

    print("\n1. Loading brand data...")
    brands = define_brand_data()
    print(f"   {len(brands)} brands loaded")

    print("\n2. Generating positioning map...")
    generate_positioning_map(brands, chart_path)

    print("\n3. Building slides...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        ("Slide  1: Title",                    lambda: build_slide_01_title(prs)),
        ("Slide  2: Retention Problem",        lambda: build_slide_02_retention(prs)),
        ("Slide  3: The Backlash",             lambda: build_slide_03_backlash(prs)),
        ("Slide  4: The Gender Gap",           lambda: build_slide_04_gender_gap(prs)),
        ("Slide  5: Positioning Map",          lambda: build_slide_05_map(prs, chart_path)),
        ("Slide  6: Optimization Quadrant",    lambda: build_slide_06_optimization(prs)),
        ("Slide  7: Lifestyle Quadrant",       lambda: build_slide_07_lifestyle(prs)),
        ("Slide  8: Anti-Establishment",       lambda: build_slide_08_anti_establishment(prs)),
        ("Slide  9: Proposed Position",        lambda: build_slide_09_proposed_position(prs)),
        ("Slide 10: Integration Moat",         lambda: build_slide_10_integration(prs)),
        ("Slide 11: Market Sizing",            lambda: build_slide_11_market(prs)),
        ("Slide 12: Berlin Context",           lambda: build_slide_12_berlin(prs)),
        ("Slide 13: Close",                    lambda: build_slide_13_close(prs)),
    ]

    for label, builder in builders:
        print(f"   {label}")
        builder()

    print(f"\n4. Saving deck...")
    prs.save(output_path)
    print(f"   Saved: {output_path}")

    print(f"\n{'=' * 45}")
    print(f"Done. {len(prs.slides)} slides generated.")
    print(f"  Deck:  {output_path}")
    print(f"  Chart: {chart_path}")


if __name__ == "__main__":
    build_deck()
